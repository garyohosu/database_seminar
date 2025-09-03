#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
新入社員向けデータベース講座PowerPointスライド生成スクリプト
初心者に優しいデータベース講座のスライドを自動生成します
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR
import os

def create_database_presentation():
    """データベース講座のPowerPointプレゼンテーションを作成"""
    
    # プレゼンテーション作成
    prs = Presentation()
    
    # スライドサイズを16:9に設定
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # カラーテーマ設定
    primary_color = RGBColor(54, 96, 146)      # ネイビーブルー
    accent_color = RGBColor(79, 129, 189)      # ライトブルー
    text_color = RGBColor(68, 68, 68)          # ダークグレー
    highlight_color = RGBColor(255, 192, 0)    # オレンジ
    
    # スライド1: タイトルスライド
    slide = prs.slides.add_slide(prs.slide_layouts[0])  # タイトルレイアウト
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "新入社員向けデータベース講座"
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.color.rgb = primary_color
    
    subtitle.text = "〜 難しくない！身近で役立つデータベース 〜\n\nACCESSで学ぶ120分の実践講座"
    subtitle.text_frame.paragraphs[0].font.size = Pt(24)
    subtitle.text_frame.paragraphs[0].font.color.rgb = accent_color
    
    # スライド2: アジェンダ
    slide = prs.slides.add_slide(prs.slide_layouts[1])  # タイトルとコンテンツ
    slide.shapes.title.text = "本日の学習内容（120分）"
    
    content = slide.placeholders[1].text_frame
    content.text = """1. 情報整理って大切ですよね（15分）
2. 表でつながる仕組み（20分）  
3. 整理整頓の3つのステップ（25分）
4. 休憩（10分）
5. ACCESS：今日使うツール（10分）
6. 実習：ACCESSで実際にやってみよう（35分）
7. プログラムから使ってみよう（10分）
8. まとめと次のステップ（15分）"""
    
    for paragraph in content.paragraphs:
        paragraph.font.size = Pt(20)
        paragraph.font.color.rgb = text_color
    
    # スライド3: セクション1開始
    slide = prs.slides.add_slide(prs.slide_layouts[2])  # セクションヘッダー
    slide.shapes.title.text = "第1部：情報整理って大切ですよね"
    slide.placeholders[1].text = "身近な困った経験から始めましょう"
    
    # スライド4: 身近な困った経験
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "こんな困った経験、ありませんか？"
    
    content = slide.placeholders[1].text_frame
    content.text = """📱 スマホの連絡先がぐちゃぐちゃ

📚 お気に入りのレシピがどこにあるか分からない

💳 いつどこで何を買ったか思い出せない

📝 メモした大事なことが見つからない"""
    
    for paragraph in content.paragraphs:
        paragraph.font.size = Pt(24)
        paragraph.font.color.rgb = text_color
    
    # スライド5: 整理の進化
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "情報整理の進化"
    
    content = slide.placeholders[1].text_frame
    content.text = """1. 手作業：ノートに書く
   → 字が読めない、なくす

2. Excel：表で管理
   → 便利だけど限界がある

3. データベース：もっと便利な整理方法！
   → 今日覚える「情報をとても上手に整理してくれる仕組み」"""
    
    for paragraph in content.paragraphs:
        paragraph.font.size = Pt(20)
        paragraph.font.color.rgb = text_color
    
    # スライド6: データベースは身近
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "データベースって身近にある"
    
    content = slide.placeholders[1].text_frame
    content.text = """実は毎日使っています！

📚 図書館の本検索システム
🛒 ネットショッピングの商品検索  
🏥 病院の診察券システム
🏢 会社の社員名簿

「難しそう」→「実は身近で便利なもの」"""
    
    for paragraph in content.paragraphs:
        paragraph.font.size = Pt(22)
        paragraph.font.color.rgb = text_color
    
    # スライド7: セクション2開始
    slide = prs.slides.add_slide(prs.slide_layouts[2])
    slide.shapes.title.text = "第2部：表でつながる仕組み"
    slide.placeholders[1].text = "学校の例で考えてみましょう"
    
    # スライド8: 学校の例（問題）
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "クラス名簿を作るとき"
    
    content = slide.placeholders[1].text_frame
    content.text = """1つの表で全部書くと...

生徒名 | 学年 | クラス | 担任の先生 | 先生の電話
田中   | 2年  | A組   | 山田先生   | 090-xxxx
佐藤   | 2年  | A組   | 山田先生   | 090-xxxx

何か変じゃない？
• 山田先生の情報が2回も書かれてる
• 先生の電話番号が変わったら全部直す？"""
    
    for paragraph in content.paragraphs:
        paragraph.font.size = Pt(18)
        paragraph.font.color.rgb = text_color
    
    # スライド9: 学校の例（解決）
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "表を分けて整理"
    
    content = slide.placeholders[1].text_frame
    content.text = """生徒の表
生徒名 | 学年 | クラス番号
田中   | 2年  | A001
佐藤   | 2年  | A001

クラスの表  
クラス番号 | クラス名 | 担任先生 | 先生電話
A001      | 2年A組   | 山田先生 | 090-xxxx

メリット
• 先生の情報は1回だけ
• 電話番号が変わっても1カ所だけ直せばOK"""
    
    for paragraph in content.paragraphs:
        paragraph.font.size = Pt(16)
        paragraph.font.color.rgb = text_color
    
    # スライド10: キーの概念
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "つなぐための印：キー"
    
    content = slide.placeholders[1].text_frame
    content.text = """つなぐための特別な印

• 学生証番号：生徒を区別するための番号（重複しない）
• クラス番号：クラスを区別するための番号

これを「キー」と呼びます
• 鍵のように、情報同士をつなぐ大切なもの
• 絶対に重複してはダメ（同じ学生証番号の人は2人いない）"""
    
    for paragraph in content.paragraphs:
        paragraph.font.size = Pt(20)
        paragraph.font.color.rgb = text_color
    
    # スライド11: なぜこの方法が良いの？
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "なぜこの方法が良いの？"
    
    content = slide.placeholders[1].text_frame
    content.text = """整理整頓のメリット

✓ 情報の重複がない → 容量の節約
✓ 更新が楽 → 1カ所直せば全体に反映
✓ 間違いが減る → 矛盾が起きにくい
✓ 検索が早い → 必要な情報をすぐ見つけられる

この整理方法を「リレーショナルデータベース」と呼びます
（「リレーション」= つながり）"""
    
    for paragraph in content.paragraphs:
        paragraph.font.size = Pt(20)
        paragraph.font.color.rgb = text_color
    
    # スライド12: セクション3開始
    slide = prs.slides.add_slide(prs.slide_layouts[2])
    slide.shapes.title.text = "第3部：整理整頓の3つのステップ"
    slide.placeholders[1].text = "お買い物リストの整理を例に"
    
    # スライド13: 問題のあるリスト
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "よくある問題：重複だらけのリスト"
    
    content = slide.placeholders[1].text_frame
    content.text = """お客様 | 住所      | 商品     | 値段 | 個数
田中   | 東京都... | りんご   | 100  | 3
田中   | 東京都... | バナナ   | 150  | 2
佐藤   | 大阪府... | りんご   | 100  | 1

問題点を見つけてみよう
• 田中さんの住所が2回書かれてる
• りんごの値段が2回書かれてる
• もし住所が変わったら？値段が変わったら？"""
    
    for paragraph in content.paragraphs:
        paragraph.font.size = Pt(18)
        paragraph.font.color.rgb = text_color
    
    # スライド14: ステップ1
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "整理のステップ1：1つのマスに1つの情報"
    
    content = slide.placeholders[1].text_frame
    content.text = """Before（ダメな例）
好きな食べ物：りんご、バナナ、みかん

After（良い例）  
好きな食べ物：りんご
好きな食べ物：バナナ
好きな食べ物：みかん

ルール：1つのマスには1つの情報だけ
これを「第1正規形」と言います（覚えなくてもOK）"""
    
    for paragraph in content.paragraphs:
        paragraph.font.size = Pt(20)
        paragraph.font.color.rgb = text_color
    
    # スライド15: ステップ2
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "整理のステップ2：関連する情報をグループ分け"
    
    content = slide.placeholders[1].text_frame
    content.text = """お客様の情報と商品の情報を分ける

お客様の表
お客様ID | お客様名 | 住所
C001     | 田中     | 東京都...
C002     | 佐藤     | 大阪府...

商品の表
商品ID | 商品名 | 値段
P001   | りんご | 100
P002   | バナナ | 150

購入の表
お客様ID | 商品ID | 個数
C001     | P001   | 3
C001     | P002   | 2
C002     | P001   | 1"""
    
    for paragraph in content.paragraphs:
        paragraph.font.size = Pt(16)
        paragraph.font.color.rgb = text_color
    
    # スライド16: ステップ3
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "整理のステップ3：間接的な関係も整理"
    
    content = slide.placeholders[1].text_frame
    content.text = """さらに詳しく分けることもあります

例：商品に「カテゴリ」がある場合
商品ID | 商品名 | カテゴリID
カテゴリID | カテゴリ名

でも今日は深く考えなくてOK！
「情報を関連するもの同士でグループ分けする」
ということが分かれば十分です"""
    
    for paragraph in content.paragraphs:
        paragraph.font.size = Pt(20)
        paragraph.font.color.rgb = text_color
    
    # スライド17: 整理整頓のメリット
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "整理整頓のメリットを実感"
    
    content = slide.placeholders[1].text_frame
    content.text = """整理前 vs 整理後

整理前の問題
• 情報の重複
• 更新の手間  
• 間違いのリスク

整理後の良さ
• 情報は1回だけ保存
• 1カ所直せばOK
• 矛盾が起きない

これが「正規化」の効果です"""
    
    for paragraph in content.paragraphs:
        paragraph.font.size = Pt(22)
        paragraph.font.color.rgb = text_color
    
    # スライド18: 休憩
    slide = prs.slides.add_slide(prs.slide_layouts[2])
    slide.shapes.title.text = "小休憩"
    slide.placeholders[1].text = "☕ ちょっと一息\n\nここまでのおさらい\n1. データベース = 情報の上手な整理方法\n2. 表を分けてつなぐ = リレーショナル\n3. 整理整頓のステップ = 正規化\n\n💬 質問タイム\n分からないことがあれば、お気軽にどうぞ！"
    
    # スライド19: セクション4開始
    slide = prs.slides.add_slide(prs.slide_layouts[2])
    slide.shapes.title.text = "第4部：ACCESS - 今日使うツール"
    slide.placeholders[1].text = "初心者にやさしい選択"
    
    # スライド20: なぜACCESS？
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "なぜACCESSを選んだの？"
    
    content = slide.placeholders[1].text_frame
    content.text = """初心者にやさしい理由

🏢 多くの会社のパソコンに入っている
🔰 操作が比較的分かりやすい  
📊 Excelからステップアップしやすい
💻 1台のパソコンで完結する

他にもいろんなデータベースがあります
MySQL、PostgreSQL、SQL Server...
でも今日はACCESSだけ覚えればOK！"""
    
    for paragraph in content.paragraphs:
        paragraph.font.size = Pt(20)
        paragraph.font.color.rgb = text_color
    
    # スライド21: 将来の選択肢
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "将来の選択肢（参考程度）"
    
    content = slide.placeholders[1].text_frame
    content.text = """規模が大きくなったら

• 小さなプロジェクト → ACCESS
• 会社のシステム → SQL Server
• Webサイト → MySQL
• 大企業 → Oracle

今は「ACCESSで十分」
覚えたことは他でも応用できます"""
    
    for paragraph in content.paragraphs:
        paragraph.font.size = Pt(22)
        paragraph.font.color.rgb = text_color
    
    # スライド22: セクション5開始
    slide = prs.slides.add_slide(prs.slide_layouts[2])
    slide.shapes.title.text = "第5部：実習 - ACCESSで実際にやってみよう"
    slide.placeholders[1].text = "お店の商品管理を作ってみましょう"
    
    # スライド23: 今日作るもの
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "今日作るもの：お店の商品管理"
    
    content = slide.placeholders[1].text_frame
    content.text = """身近な例：小さなお店

• 商品（商品名、値段、在庫数）
• お客様（名前、電話番号）
• 売上（いつ、誰が、何を、いくつ買った）

段階的に進めます
1. データを見る
2. データを追加
3. データを変更
4. データを削除"""
    
    for paragraph in content.paragraphs:
        paragraph.font.size = Pt(20)
        paragraph.font.color.rgb = text_color
    
    # スライド24: SQLの紹介
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "魔法の言葉：SQL"
    
    content = slide.placeholders[1].text_frame
    content.text = """SQL = データベースに指示を出す言葉

• 英語に似ているので分かりやすい
• SELECT = 選ぶ
• FROM = ～から
• WHERE = ～という条件で

最初の魔法
SELECT * FROM 商品;

意味：「商品テーブルから、全部（*）を選んで表示」"""
    
    for paragraph in content.paragraphs:
        paragraph.font.size = Pt(18)
        paragraph.font.color.rgb = text_color
    
    # スライド25: 実習1 - SELECT
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "実習1：データを見てみよう"
    
    content = slide.placeholders[1].text_frame
    content.text = """みんなで一緒にやりましょう

-- 商品を全部見る
SELECT * FROM 商品;

-- 商品名と値段だけ見る
SELECT 商品名, 値段 FROM 商品;

-- 100円以上の商品だけ見る  
SELECT * FROM 商品 WHERE 値段 >= 100;

ポイント
• 「*」は「全部」という意味
• 「--」はメモ（コメント）
• セミコロン「;」で終わる"""
    
    for paragraph in content.paragraphs:
        paragraph.font.size = Pt(16)
        paragraph.font.color.rgb = text_color
    
    # スライド26: 実習2 - INSERT
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "実習2：データを追加してみよう"
    
    content = slide.placeholders[1].text_frame
    content.text = """新しい商品を追加

-- 新商品を追加
INSERT INTO 商品 (商品名, 値段, 在庫数)
VALUES ('消しゴム', 80, 20);

意味
• INSERT INTO = ～に挿入
• VALUES = 値は～

やってみよう
お好きな商品を1つ追加してください！"""
    
    for paragraph in content.paragraphs:
        paragraph.font.size = Pt(18)
        paragraph.font.color.rgb = text_color
    
    # スライド27: 実習3 - UPDATE
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "実習3：データを変更してみよう"
    
    content = slide.placeholders[1].text_frame
    content.text = """値段を変更する

-- 消しゴムの値段を変更
UPDATE 商品
SET 値段 = 90
WHERE 商品名 = '消しゴム';

⚠️ 大事な注意
WHEREを忘れると、全部の商品の値段が変わっちゃいます！"""
    
    for paragraph in content.paragraphs:
        paragraph.font.size = Pt(20)
        paragraph.font.color.rgb = text_color
    
    # スライド28: 実習4 - DELETE
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "実習4：データを削除してみよう"
    
    content = slide.placeholders[1].text_frame
    content.text = """在庫がない商品を削除

-- 在庫0の商品を削除
DELETE FROM 商品
WHERE 在庫数 = 0;

⚠️ さらに重要な注意
WHEREを忘れると、全部の商品が消えちゃいます！"""
    
    for paragraph in content.paragraphs:
        paragraph.font.size = Pt(20)
        paragraph.font.color.rgb = text_color
    
    # スライド29: 実習5 - JOIN
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "実習5：2つの表をつなげてみよう"
    
    content = slide.placeholders[1].text_frame
    content.text = """誰が何を買ったか見る

-- お客様と売上をつなげて見る
SELECT お客様.名前, 商品.商品名, 売上.数量
FROM お客様, 商品, 売上
WHERE お客様.お客様ID = 売上.お客様ID
  AND 商品.商品ID = 売上.商品ID;

少し複雑ですが
「表同士をつなげて、欲しい情報を取り出している」
ということが分かればOK！"""
    
    for paragraph in content.paragraphs:
        paragraph.font.size = Pt(16)
        paragraph.font.color.rgb = text_color
    
    # スライド30: セクション6開始
    slide = prs.slides.add_slide(prs.slide_layouts[2])
    slide.shapes.title.text = "第6部：プログラムから使ってみよう"
    slide.placeholders[1].text = "C++Builderからデータベースを操作"
    
    # スライド31: 接続方法
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "C++Builderから接続"
    
    content = slide.placeholders[1].text_frame
    content.text = """プログラムからもデータベースを使えます

接続の準備（コピペでOK）
// データベースに接続する準備
String connectionString = 
  "Provider=Microsoft.ACE.OLEDB.12.0;"  // ACCESSを使うよという宣言
  "Data Source=C:\\\\sample.accdb;";      // ファイルの場所

// 接続実行
ADOConnection1->ConnectionString = connectionString;
ADOConnection1->Connected = true;  // 接続開始"""
    
    for paragraph in content.paragraphs:
        paragraph.font.size = Pt(14)
        paragraph.font.color.rgb = text_color
    
    # スライド32: データ読み取り
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "データを取得するプログラム"
    
    content = slide.placeholders[1].text_frame
    content.text = """try {
  // SQLを準備
  ADOQuery1->SQL->Clear();                    // 前のSQLをクリア
  ADOQuery1->SQL->Add("SELECT * FROM 商品");  // SQLを追加
  ADOQuery1->Open();                          // 実行
  
  // 結果を1行ずつ見る
  while (!ADOQuery1->Eof) {  // 終わりまで繰り返し
    String name = ADOQuery1->FieldByName("商品名")->AsString;
    int price = ADOQuery1->FieldByName("値段")->AsInteger;
    
    ShowMessage(name + "の値段は" + IntToStr(price) + "円");
    ADOQuery1->Next();  // 次の行へ
  }
} catch (Exception &e) {
  ShowMessage("エラー: " + e.Message);
}"""
    
    for paragraph in content.paragraphs:
        paragraph.font.size = Pt(12)
        paragraph.font.color.rgb = text_color
    
    # スライド33: データ追加プログラム
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "データを追加するプログラム"
    
    content = slide.placeholders[1].text_frame
    content.text = """try {
  // パラメータ付きSQL（安全な方法）
  ADOQuery1->SQL->Clear();
  ADOQuery1->SQL->Add("INSERT INTO 商品 (商品名, 値段, 在庫数) "
                      "VALUES (:name, :price, :stock)");
  
  // 値を設定
  ADOQuery1->Parameters->ParamByName("name")->Value = "新商品";
  ADOQuery1->Parameters->ParamByName("price")->Value = 200;
  ADOQuery1->Parameters->ParamByName("stock")->Value = 10;
  
  ADOQuery1->ExecSQL();  // 実行
  ShowMessage("商品を追加しました！");
  
} catch (Exception &e) {
  ShowMessage("追加に失敗: " + e.Message);
}

重要ポイント
• try-catchでエラー対策必須
• パラメータ（:name等）で安全に値を渡す"""
    
    for paragraph in content.paragraphs:
        paragraph.font.size = Pt(12)
        paragraph.font.color.rgb = text_color
    
    # スライド34: セクション7開始
    slide = prs.slides.add_slide(prs.slide_layouts[2])
    slide.shapes.title.text = "第7部：まとめと次のステップ"
    slide.placeholders[1].text = "今日の成果を確認しましょう"
    
    # スライド35: 今日覚えたこと
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "今日覚えたこと"
    
    content = slide.placeholders[1].text_frame
    content.text = """✅ データベースの基本
• 情報を整理する便利な仕組み
• 身近なところで使われている
• 難しくない、便利な道具

✅ 表の整理方法
• 関連する情報はまとめる
• 重複は避ける
• キーでつなぐ

✅ SQLの基本
• SELECT：見る  • INSERT：追加
• UPDATE：変更 • DELETE：削除

✅ プログラムからの利用
• 接続文字列で接続  • エラー対策が重要"""
    
    for paragraph in content.paragraphs:
        paragraph.font.size = Pt(16)
        paragraph.font.color.rgb = text_color
    
    # スライド36: 忘れても大丈夫
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "忘れても大丈夫なこと"
    
    content = slide.placeholders[1].text_frame
    content.text = """細かい文法 → 必要な時に調べればOK
複雑な理論 → 実際に使いながら覚える
専門用語 → 概念が分かっていれば十分

大切なのは「考え方」
• 情報を整理する発想
• データの関連性を意識する
• 段階的に学習する姿勢"""
    
    for paragraph in content.paragraphs:
        paragraph.font.size = Pt(20)
        paragraph.font.color.rgb = text_color
    
    # スライド37: 次に学ぶと良いこと
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "次に学ぶと良いこと（将来の参考）"
    
    content = slide.placeholders[1].text_frame
    content.text = """もっとSQLを使いたい場合
• JOIN（表の結合）の詳しい使い方
• GROUP BY（集計）
• 関数（SUM、COUNT等）

もっとACCESSを使いたい場合
• フォームの作成
• レポートの作成
• マクロの活用

他のデータベースに挑戦
• MySQL（Web開発）
• SQL Server（企業システム）"""
    
    for paragraph in content.paragraphs:
        paragraph.font.size = Pt(18)
        paragraph.font.color.rgb = text_color
    
    # スライド38: よくある質問
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "よくある質問"
    
    content = slide.placeholders[1].text_frame
    content.text = """Q: SQLを全部覚える必要がある？
A: いいえ。基本の4つができれば十分。必要に応じて調べましょう。

Q: プログラムが難しそう...
A: 最初はコピペから始めてOK。少しずつ理解していけば大丈夫。

Q: 実際の業務でどう使う？
A: 顧客管理、在庫管理、売上分析など。小さなところから始めて。

Q: データが消えちゃったらどうする？
A: バックアップが重要！定期的にファイルをコピーしておきましょう。"""
    
    for paragraph in content.paragraphs:
        paragraph.font.size = Pt(16)
        paragraph.font.color.rgb = text_color
    
    # スライド39: 最後のメッセージ
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "最後に"
    
    content = slide.placeholders[1].text_frame
    content.text = """🎉 お疲れさまでした！

今日の成果
• データベースの世界に第一歩
• SQLで実際にデータを操作
• プログラムからの操作も体験

これからも
• 分からないことがあっても焦らない
• 小さなプロジェクトから始める
• 困った時はいつでも相談してください

データベースは怖くない、とても便利な道具です！"""
    
    for paragraph in content.paragraphs:
        paragraph.font.size = Pt(20)
        paragraph.font.color.rgb = text_color
    
    return prs

def main():
    """メイン実行関数"""
    print("新入社員向けデータベース講座のPowerPointスライドを生成中...")
    
    try:
        # プレゼンテーション作成
        presentation = create_database_presentation()
        
        # ファイル保存
        output_file = "新入社員向けデータベース講座.pptx"
        presentation.save(output_file)
        
        print(f"スライド生成完了: {output_file}")
        print(f"総スライド数: {len(presentation.slides)}枚")
        print("初心者に優しいデータベース講座スライドが完成しました！")
        
        # ファイルの絶対パスを表示
        abs_path = os.path.abspath(output_file)
        print(f"保存場所: {abs_path}")
        
    except Exception as e:
        print(f"エラーが発生しました: {str(e)}")
        return False
    
    return True

if __name__ == "__main__":
    main()